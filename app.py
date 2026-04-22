"""
================================================================================
Ultimate AIO Epic Slide Generator
================================================================================
A Full-Stack Streamlit application that automates 100% of professional
presentation creation by integrating:
    • Gemini LLM API    → Text structuring & slide layout
    • DALL-E 3 API      → Cinematic background image generation
    • python-pptx       → PowerPoint rendering
    • Pillow (PIL)      → Image overlay & compositing

All file operations are performed in-memory (RAM) via io.BytesIO for
Cloud-native deployment. No temporary disk writes.

Author: AI Engineer
Date: 2026-04-22
================================================================================
"""

import json
import asyncio
import base64
import re
from io import BytesIO
from typing import List, Dict, Any, Optional, Tuple

import streamlit as st
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# LLM & Image Generation APIs
import google.generativeai as genai
from openai import AsyncOpenAI


# =============================================================================
# INTERNATIONALIZATION (I18N)
# =============================================================================

I18N = {
    "en": {
        "page_title": "Ultimate AIO Epic Slide Generator",
        "header_title": "🎬 Ultimate AIO Epic Slide Generator",
        "header_subtitle": "Paste raw text. Let AI design cinematic slides.",
        "sidebar_api": "🔑 API Configuration",
        "gemini_key": "Gemini API Key",
        "openai_key": "OpenAI API Key",
        "sidebar_settings": "⚙️ Presentation Settings",
        "lang_label": "Output Language",
        "theme_label": "Visual Theme",
        "slide_count_label": "Number of Slides",
        "main_input_label": "📄 Raw Text Input",
        "main_input_help": "Paste any raw content, article, or outline. The AI will auto-structure it into slides.",
        "generate_btn": "🚀 Auto-Generate Presentation",
        "progress_analyzing": "🧠 Step 1/3: Analyzing text with Gemini LLM...",
        "progress_images": "🎨 Step 2/3: Generating background images with DALL-E 3...",
        "progress_rendering": "📦 Step 3/3: Rendering final PPTX...",
        "success_msg": "🎉 Your epic presentation is ready!",
        "download_btn": "📥 Download .pptx",
        "error_no_text": "⚠️ Please enter some raw text before generating.",
        "error_no_gemini_key": "⚠️ Please provide your Gemini API Key.",
        "error_no_openai_key": "⚠️ Please provide your OpenAI API Key.",
        "error_gemini": "❌ Gemini API Error: {}",
        "error_openai": "❌ DALL-E 3 API Error: {}",
        "error_generic": "❌ Unexpected Error: {}",
        "warning_image_fallback": "⚠️ Some images failed. Using fallback backgrounds.",
        "status_ready": "Ready",
        "status_generating": "Generating...",
        "footer": "Ultimate AIO Epic Slide Generator | Built with Streamlit + Gemini + DALL-E 3",
    },
    "vi": {
        "page_title": "Ultimate AIO Epic Slide Generator",
        "header_title": "🎬 Ultimate AIO Epic Slide Generator",
        "header_subtitle": "Dán văn bản thô. Để AI thiết kế slide điện ảnh.",
        "sidebar_api": "🔑 Cấu hình API",
        "gemini_key": "Gemini API Key",
        "openai_key": "OpenAI API Key",
        "sidebar_settings": "⚙️ Cài đặt bài thuyết trình",
        "lang_label": "Ngôn ngữ đầu ra",
        "theme_label": "Chủ đề trực quan",
        "slide_count_label": "Số lượng Slide",
        "main_input_label": "📄 Văn bản thô",
        "main_input_help": "Dán bất kỳ nội dung, bài viết hoặc dàn ý nào. AI sẽ tự động cấu trúc thành các slide.",
        "generate_btn": "🚀 Tự động tạo bài thuyết trình",
        "progress_analyzing": "🧠 Bước 1/3: Đang phân tích văn bản bằng Gemini LLM...",
        "progress_images": "🎨 Bước 2/3: Đang sinh ảnh nền bằng DALL-E 3...",
        "progress_rendering": "📦 Bước 3/3: Đang render file PPTX...",
        "success_msg": "🎉 Bài thuyết trình đã sẵn sàng!",
        "download_btn": "📥 Tải xuống .pptx",
        "error_no_text": "⚠️ Vui lòng nhập văn bản trước khi tạo.",
        "error_no_gemini_key": "⚠️ Vui lòng cung cấp Gemini API Key.",
        "error_no_openai_key": "⚠️ Vui lòng cung cấp OpenAI API Key.",
        "error_gemini": "❌ Lỗi API Gemini: {}",
        "error_openai": "❌ Lỗi API DALL-E 3: {}",
        "error_generic": "❌ Lỗi không mong muốn: {}",
        "warning_image_fallback": "⚠️ Một số ảnh sinh thất bại. Sử dụng ảnh nền dự phòng.",
        "status_ready": "Sẵn sàng",
        "status_generating": "Đang tạo...",
        "footer": "Ultimate AIO Epic Slide Generator | Xây dựng bởi Streamlit + Gemini + DALL-E 3",
    },
}


def get_text(key: str, lang: str = "en") -> str:
    """Retrieve a translated string by key for the selected language."""
    return I18N.get(lang, I18N["en"]).get(key, key)


# =============================================================================
# THEME ENGINE
# =============================================================================
# Each theme drives: typography, colors, positioning, and DALL-E prompt keywords.

THEMES: Dict[str, Dict[str, Any]] = {
    "Epic Sports": {
        "name": "Epic Sports",
        "title_font": "Impact",
        "bullet_font": "Arial",
        "title_color": RGBColor(255, 215, 0),        # Gold
        "bullet_color": RGBColor(255, 255, 255),     # White
        "overlay_rgb": (0, 0, 0),
        "overlay_opacity": 0.60,
        "dalle_keywords": "cinematic lighting, hyper-realistic, epic atmosphere, dramatic shadows, sports arena, high contrast",
        "title_pos": {"left": 0.5, "top": 1.5, "width": 7.0, "height": 2.0},
        "bullet_pos": {"left": 0.5, "top": 4.0, "width": 7.0, "height": 4.0},
        "overlay_width_ratio": 0.50,
    },
    "Corporate": {
        "name": "Corporate",
        "title_font": "Arial",
        "bullet_font": "Arial",
        "title_color": RGBColor(255, 255, 255),       # White
        "bullet_color": RGBColor(200, 220, 240),      # Light steel blue
        "overlay_rgb": (25, 40, 65),
        "overlay_opacity": 0.75,
        "dalle_keywords": "minimal corporate background, clean lines, professional lighting, navy blue tones, modern office, subtle gradient",
        "title_pos": {"left": 0.75, "top": 1.5, "width": 6.5, "height": 2.0},
        "bullet_pos": {"left": 0.75, "top": 4.0, "width": 6.5, "height": 4.0},
        "overlay_width_ratio": 0.50,
    },
    "Cyberpunk Neon": {
        "name": "Cyberpunk Neon",
        "title_font": "Consolas",
        "bullet_font": "Consolas",
        "title_color": RGBColor(0, 255, 255),         # Cyan
        "bullet_color": RGBColor(255, 0, 255),        # Magenta
        "overlay_rgb": (10, 0, 20),
        "overlay_opacity": 0.70,
        "dalle_keywords": "cyberpunk neon city, glowing lights, futuristic, dark background, purple and cyan neon, sci-fi",
        "title_pos": {"left": 0.5, "top": 1.2, "width": 7.0, "height": 2.0},
        "bullet_pos": {"left": 0.5, "top": 3.8, "width": 7.0, "height": 4.5},
        "overlay_width_ratio": 0.55,
    },
}


# =============================================================================
# GLOBAL CONSTANTS
# =============================================================================

SLIDE_WIDTH = Inches(16)
SLIDE_HEIGHT = Inches(9)
DPI = 96  # pixels per inch for PIL rendering

# DALL-E 3 landscape size closest to 16:9
DALLE_SIZE = "1792x1024"


# =============================================================================
# HELPER FUNCTIONS — IMAGE & OVERLAY GENERATION
# =============================================================================

def create_overlay_image(theme: Dict[str, Any]) -> BytesIO:
    """
    Generate a left-half overlay PNG using Pillow.
    The overlay covers the left portion of the slide (configurable by theme)
    with theme-specific color and opacity.

    Args:
        theme: The selected theme dictionary.

    Returns:
        BytesIO buffer containing a PNG image with alpha transparency.
    """
    # Calculate overlay dimensions in pixels
    width_inches = 16 * theme["overlay_width_ratio"]
    height_inches = 9
    width_px = int(width_inches * DPI)
    height_px = int(height_inches * DPI)

    r, g, b = theme["overlay_rgb"]
    alpha = int(255 * theme["overlay_opacity"])

    overlay_rgba = PILImage.new("RGBA", (width_px, height_px), (r, g, b, alpha))

    img_buffer = BytesIO()
    overlay_rgba.save(img_buffer, format="PNG")
    img_buffer.seek(0)
    return img_buffer


def create_fallback_gradient(theme: Dict[str, Any]) -> BytesIO:
    """
    Create a fallback gradient image if DALL-E generation fails.
    Produces a smooth vertical gradient based on the theme's overlay color.

    Args:
        theme: The selected theme dictionary.

    Returns:
        BytesIO buffer containing a PNG image.
    """
    width_px = int(16 * DPI)
    height_px = int(9 * DPI)

    base_r, base_g, base_b = theme["overlay_rgb"]

    # Create a vertical gradient from dark to slightly lighter
    img = PILImage.new("RGB", (width_px, height_px))
    pixels = img.load()

    for y in range(height_px):
        ratio = y / height_px
        # Interpolate from base color to a slightly lighter variant
        r = int(base_r + (60 * ratio))
        g = int(base_g + (60 * ratio))
        b = int(base_b + (60 * ratio))
        for x in range(width_px):
            pixels[x, y] = (r, g, b)

    buffer = BytesIO()
    img.save(buffer, format="PNG")
    buffer.seek(0)
    return buffer


# =============================================================================
# HELPER FUNCTIONS — PPTX SLIDE ASSEMBLY
# =============================================================================

def add_background_image(slide, image_bytes: BytesIO) -> None:
    """
    Add a full-bleed background image to a slide.
    The image is stretched to cover the entire 16x9 slide.
    """
    slide.shapes.add_picture(
        image_bytes,
        Inches(0),
        Inches(0),
        width=SLIDE_WIDTH,
        height=SLIDE_HEIGHT
    )


def add_overlay(slide, overlay_bytes: BytesIO, theme: Dict[str, Any]) -> None:
    """
    Add the semi-transparent overlay to the left portion of the slide.
    Position is always anchored at top-left.
    """
    overlay_width = Inches(16 * theme["overlay_width_ratio"])
    overlay_height = Inches(9)

    slide.shapes.add_picture(
        overlay_bytes,
        Inches(0),
        Inches(0),
        width=overlay_width,
        height=overlay_height
    )


def add_title_textbox(slide, title_text: str, theme: Dict[str, Any]) -> None:
    """
    Add a styled title textbox onto the slide.
    Text is rendered in UPPERCASE for epic visual impact.
    """
    pos = theme["title_pos"]
    textbox = slide.shapes.add_textbox(
        Inches(pos["left"]),
        Inches(pos["top"]),
        Inches(pos["width"]),
        Inches(pos["height"])
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT

    run = paragraph.add_run()
    run.text = title_text.upper()
    run.font.name = theme["title_font"]
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = theme["title_color"]


def add_bullet_textbox(slide, bullets: List[str], theme: Dict[str, Any]) -> None:
    """
    Add styled bullet points onto the slide.
    Handles UTF-8 Vietnamese text correctly via python-pptx.
    """
    pos = theme["bullet_pos"]
    textbox = slide.shapes.add_textbox(
        Inches(pos["left"]),
        Inches(pos["top"]),
        Inches(pos["width"]),
        Inches(pos["height"])
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for idx, bullet_text in enumerate(bullets):
        if idx == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()

        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.level = 0

        run = paragraph.add_run()
        run.text = str(bullet_text)
        run.font.name = theme["bullet_font"]
        run.font.size = Pt(24)
        run.font.color.rgb = theme["bullet_color"]


def build_slide(
    slide,
    slide_data: Dict[str, Any],
    image_bytes: BytesIO,
    overlay_bytes: BytesIO,
    theme: Dict[str, Any]
) -> None:
    """
    Assemble a single slide with:
        1. Full-bleed background image
        2. Semi-transparent overlay on the left
        3. Title text
        4. Bullet points
    """
    # Layer 1: Background (full-bleed)
    add_background_image(slide, image_bytes)

    # Layer 2: Overlay (left half, semi-transparent)
    add_overlay(slide, overlay_bytes, theme)

    # Layer 3: Title
    title = slide_data.get("title", "")
    if title:
        add_title_textbox(slide, title, theme)

    # Layer 4: Bullets
    bullets = slide_data.get("bullets", [])
    if bullets:
        add_bullet_textbox(slide, bullets, theme)


def build_presentation(
    slides_data: List[Dict[str, Any]],
    image_buffers: List[BytesIO],
    theme: Dict[str, Any]
) -> BytesIO:
    """
    Build the complete PowerPoint presentation in memory.

    Args:
        slides_data: List of slide dictionaries (title, bullets, etc.)
        image_buffers: List of BytesIO images (one per slide)
        theme: The selected theme dictionary

    Returns:
        BytesIO buffer containing the final .pptx file
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # Pre-generate overlay once (same for all slides in this theme)
    overlay_buffer = create_overlay_image(theme)

    for slide_data, img_buffer in zip(slides_data, image_buffers):
        slide = prs.slides.add_slide(blank_layout)
        build_slide(slide, slide_data, img_buffer, overlay_buffer, theme)

    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer


# =============================================================================
# API INTEGRATION — GEMINI LLM
# =============================================================================

def build_gemini_prompt(raw_text: str, lang: str, num_slides: int, theme: Dict[str, Any]) -> str:
    """
    Construct a strict system prompt for Gemini to return structured JSON.
    The prompt instructs the model to analyze raw text and emit slide data.
    """
    lang_instruction = (
        "Vietnamese (Tiếng Việt)"
        if lang == "vi"
        else "English"
    )

    prompt = f"""You are an expert presentation designer and copywriter.

TASK:
Analyze the raw text provided below and structure it into exactly {num_slides} slides for a professional presentation.

OUTPUT FORMAT:
Return ONLY a valid JSON object (no markdown, no code blocks, no extra commentary). The JSON must follow this exact schema:

{{
  "slides": [
    {{
      "title": "A compelling, concise title for this slide",
      "bullets": [
        "First key point (1 short sentence)",
        "Second key point (1 short sentence)"
      ],
      "image_generation_prompt": "A detailed English prompt for DALL-E 3 to generate a cinematic background image that matches the slide topic and the theme: {theme['name']}. Do NOT include text, words, or letters in the image."
    }}
  ]
}}

RULES:
1. Output language for "title" and "bullets" must be: {lang_instruction}.
2. "image_generation_prompt" must ALWAYS be in English.
3. Each slide should have 2 to 4 bullet points.
4. Titles should be impactful and professional.
5. The image_generation_prompt should describe a visual scene related to the slide content, incorporating these style keywords: {theme['dalle_keywords']}.
6. Ensure the JSON is syntactically valid. Escape quotes properly.

RAW TEXT TO ANALYZE:
---
{raw_text}
---
"""
    return prompt


def call_gemini_api(
    raw_text: str,
    gemini_key: str,
    lang: str,
    num_slides: int,
    theme: Dict[str, Any]
) -> List[Dict[str, Any]]:
    """
    Synchronously call the Gemini API to structure raw text into slide data.

    Args:
        raw_text: The user's raw text input.
        gemini_key: Gemini API key.
        lang: Output language code.
        num_slides: Target number of slides.
        theme: Selected theme dictionary.

    Returns:
        List of slide dictionaries.

    Raises:
        Exception: If the API call fails or JSON parsing fails.
    """
    # Configure Gemini client
    genai.configure(api_key=gemini_key)
    model = genai.GenerativeModel("gemini-1.5-flash")

    prompt = build_gemini_prompt(raw_text, lang, num_slides, theme)

    try:
        response = model.generate_content(prompt)
    except Exception as e:
        raise RuntimeError(f"Gemini API request failed: {e}")

    if not response.text:
        raise RuntimeError("Gemini returned an empty response.")

    # Extract JSON from response (handle markdown code blocks robustly)
    raw_response = response.text.strip()

    # Try to extract from ```json ... ``` or ``` ... ```
    json_match = re.search(r"```(?:json)?\s*([\s\S]*?)\s*```", raw_response)
    if json_match:
        json_str = json_match.group(1).strip()
    else:
        json_str = raw_response

    # Fallback: find the outermost JSON object
    start_idx = json_str.find("{")
    end_idx = json_str.rfind("}")
    if start_idx == -1 or end_idx == -1 or end_idx <= start_idx:
        raise RuntimeError(f"Could not locate valid JSON in Gemini response. Raw: {raw_response[:500]}")

    json_str = json_str[start_idx:end_idx + 1]

    try:
        parsed = json.loads(json_str)
    except json.JSONDecodeError as e:
        raise RuntimeError(f"JSON decode error: {e}. Content: {json_str[:500]}")

    if "slides" not in parsed or not isinstance(parsed["slides"], list):
        raise RuntimeError("Parsed JSON does not contain a 'slides' array.")

    return parsed["slides"]


# =============================================================================
# API INTEGRATION — DALL-E 3 (ASYNC)
# =============================================================================

async def generate_image_async(
    client: AsyncOpenAI,
    prompt: str,
    size: str = DALLE_SIZE
) -> BytesIO:
    """
    Asynchronously call the OpenAI DALL-E 3 API to generate an image.

    Args:
        client: Initialized AsyncOpenAI client.
        prompt: The image generation prompt (English).
        size: Image dimensions. Default is 1792x1024 (landscape, near 16:9).

    Returns:
        BytesIO buffer containing the generated PNG image.

    Raises:
        Exception: If the API call fails.
    """
    response = await client.images.generate(
        model="dall-e-3",
        prompt=prompt,
        size=size,
        quality="standard",
        n=1,
        response_format="b64_json"
    )

    if not response.data or not response.data[0].b64_json:
        raise RuntimeError("DALL-E 3 returned empty image data.")

    image_data = base64.b64decode(response.data[0].b64_json)
    return BytesIO(image_data)


async def generate_all_images(
    openai_key: str,
    slides_data: List[Dict[str, Any]],
    theme: Dict[str, Any],
    max_concurrency: int = 3
) -> Tuple[List[BytesIO], bool]:
    """
    Concurrently generate background images for all slides using DALL-E 3.
    Uses an asyncio.Semaphore to limit concurrent API calls (rate-limit safety).

    Args:
        openai_key: OpenAI API key.
        slides_data: List of slide dictionaries containing image_generation_prompt.
        theme: Selected theme dictionary (for fallback gradients).
        max_concurrency: Maximum simultaneous DALL-E requests.

    Returns:
        Tuple of (list of image BytesIO buffers, bool indicating if any fallback was used).
    """
    client = AsyncOpenAI(api_key=openai_key, timeout=60.0)
    semaphore = asyncio.Semaphore(max_concurrency)

    async def _generate_one(slide_data: Dict[str, Any]) -> BytesIO:
        """Inner coroutine with semaphore-guarded API call + fallback."""
        base_prompt = slide_data.get("image_generation_prompt", "abstract cinematic background")
        # Enforce no-text policy and append theme keywords
        final_prompt = (
            f"{base_prompt}. {theme['dalle_keywords']}. "
            "High quality, photorealistic, NO text, NO letters, NO words, NO watermarks."
        )

        async with semaphore:
            try:
                # Small internal delay to further soften rate-limit pressure
                await asyncio.sleep(0.5)
                return await generate_image_async(client, final_prompt)
            except Exception as e:
                # Log error to Streamlit console (will be caught and surfaced by caller)
                raise RuntimeError(f"Image generation failed: {e}")

    tasks = [_generate_one(sd) for sd in slides_data]
    results = await asyncio.gather(*tasks, return_exceptions=True)

    image_buffers: List[BytesIO] = []
    any_fallback = False

    for idx, result in enumerate(results):
        if isinstance(result, Exception):
            # Fallback to gradient
            any_fallback = True
            fallback = create_fallback_gradient(theme)
            image_buffers.append(fallback)
        else:
            image_buffers.append(result)

    return image_buffers, any_fallback


# =============================================================================
# ORCHESTRATION PIPELINE
# =============================================================================

def run_generation_pipeline(
    raw_text: str,
    gemini_key: str,
    openai_key: str,
    lang: str,
    theme_name: str,
    num_slides: int,
    progress_bar,
    status_text
) -> BytesIO:
    """
    Execute the full automation pipeline:
        Step 1 → Gemini LLM text structuring
        Step 2 → Async DALL-E 3 image generation
        Step 3 → python-pptx rendering

    Args:
        raw_text: User's raw text input.
        gemini_key: Gemini API key.
        openai_key: OpenAI API key.
        lang: Selected output language.
        theme_name: Selected theme name.
        num_slides: Target number of slides.
        progress_bar: Streamlit progress bar widget.
        status_text: Streamlit empty container for status messages.

    Returns:
        BytesIO buffer containing the final .pptx file.
    """
    theme = THEMES[theme_name]

    # -------------------------------------------------------------------------
    # STEP 1: LLM Processing (Gemini)
    # -------------------------------------------------------------------------
    status_text.text(get_text("progress_analyzing", lang))
    progress_bar.progress(10)

    slides_data = call_gemini_api(
        raw_text=raw_text,
        gemini_key=gemini_key,
        lang=lang,
        num_slides=num_slides,
        theme=theme
    )

    # Ensure we don't exceed the requested number of slides
    slides_data = slides_data[:num_slides]

    progress_bar.progress(35)

    # -------------------------------------------------------------------------
    # STEP 2: Async Image Generation (DALL-E 3)
    # -------------------------------------------------------------------------
    status_text.text(get_text("progress_images", lang))
    progress_bar.progress(40)

    image_buffers, any_fallback = asyncio.run(
        generate_all_images(
            openai_key=openai_key,
            slides_data=slides_data,
            theme=theme,
            max_concurrency=3
        )
    )

    if any_fallback:
        status_text.warning(get_text("warning_image_fallback", lang))

    progress_bar.progress(75)

    # -------------------------------------------------------------------------
    # STEP 3: Render PPTX
    # -------------------------------------------------------------------------
    status_text.text(get_text("progress_rendering", lang))
    progress_bar.progress(80)

    pptx_buffer = build_presentation(
        slides_data=slides_data,
        image_buffers=image_buffers,
        theme=theme
    )

    progress_bar.progress(100)
    status_text.text(get_text("success_msg", lang))

    return pptx_buffer


# =============================================================================
# STREAMLIT UI
# =============================================================================

def render_sidebar(lang: str) -> Tuple[str, str, str, str, int]:
    """
    Render the Streamlit sidebar and collect user inputs.

    Returns:
        Tuple of (gemini_key, openai_key, selected_lang, selected_theme, num_slides)
    """
    st.sidebar.markdown(f"### {get_text('sidebar_api', lang)}")

    gemini_key = st.sidebar.text_input(
        label=get_text("gemini_key", lang),
        type="password",
        placeholder="AIzaSy...",
        help="Get your key at https://aistudio.google.com/app/apikey"
    )

    openai_key = st.sidebar.text_input(
        label=get_text("openai_key", lang),
        type="password",
        placeholder="sk-...",
        help="Get your key at https://platform.openai.com/api-keys"
    )

    st.sidebar.markdown("---")
    st.sidebar.markdown(f"### {get_text('sidebar_settings', lang)}")

    selected_lang = st.sidebar.selectbox(
        label=get_text("lang_label", lang),
        options=["en", "vi"],
        format_func=lambda x: "🇺🇸 English" if x == "en" else "🇻🇳 Tiếng Việt",
        index=0 if lang == "en" else 1
    )

    selected_theme = st.sidebar.selectbox(
        label=get_text("theme_label", lang),
        options=list(THEMES.keys()),
        index=0
    )

    num_slides = st.sidebar.slider(
        label=get_text("slide_count_label", lang),
        min_value=3,
        max_value=15,
        value=5,
        step=1
    )

    st.sidebar.markdown("---")
    st.sidebar.info(
        "💡 **Tip:**\n\n"
        "• Gemini structures your text\n"
        "• DALL-E 3 creates backgrounds\n"
        "• Everything stays in RAM\n"
        "• No files saved to disk"
    )

    return gemini_key, openai_key, selected_lang, selected_theme, num_slides


def main() -> None:
    """Main entry point for the Streamlit application."""
    st.set_page_config(
        page_title="Ultimate AIO Epic Slide Generator",
        page_icon="🎬",
        layout="wide",
        initial_sidebar_state="expanded"
    )

    # Custom CSS for cinematic feel
    st.markdown("""
    <style>
    .main-header {
        font-size: 3rem;
        font-weight: 900;
        background: linear-gradient(90deg, #FF4B4B, #FFD700);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        text-align: center;
        margin-bottom: 0.3rem;
    }
    .sub-header {
        font-size: 1.15rem;
        color: #9ca3af;
        text-align: center;
        margin-bottom: 2rem;
    }
    .stTextArea textarea {
        font-family: 'Segoe UI', sans-serif;
        font-size: 1rem;
    }
    div[data-testid="stDownloadButton"] button {
        background-color: #10b981;
        color: white;
        font-weight: bold;
        border-radius: 8px;
        padding: 0.7rem 1.4rem;
    }
    div[data-testid="stDownloadButton"] button:hover {
        background-color: #059669;
    }
    .footer {
        text-align: center;
        color: #6b7280;
        margin-top: 3rem;
        font-size: 0.85rem;
    }
    </style>
    """, unsafe_allow_html=True)

    # -------------------------------------------------------------------------
    # Language Selection (top-level, independent of sidebar return)
    # -------------------------------------------------------------------------
    # We use session state to persist language across reruns until sidebar overrides
    if "ui_lang" not in st.session_state:
        st.session_state.ui_lang = "en"

    lang = st.session_state.ui_lang

    # Header
    st.markdown(f'<div class="main-header">{get_text("header_title", lang)}</div>', unsafe_allow_html=True)
    st.markdown(f'<div class="sub-header">{get_text("header_subtitle", lang)}</div>', unsafe_allow_html=True)

    # -------------------------------------------------------------------------
    # Sidebar Inputs
    # -------------------------------------------------------------------------
    gemini_key, openai_key, selected_lang, selected_theme, num_slides = render_sidebar(lang)

    # Sync language if changed in sidebar
    if selected_lang != st.session_state.ui_lang:
        st.session_state.ui_lang = selected_lang
        st.rerun()

    # -------------------------------------------------------------------------
    # Main Input Area
    # -------------------------------------------------------------------------
    raw_text = st.text_area(
        label=get_text("main_input_label", lang),
        height=280,
        placeholder=get_text("main_input_help", lang),
        help=get_text("main_input_help", lang)
    )

    # Progress indicators (initially empty)
    progress_bar = st.progress(0)
    status_text = st.empty()

    # -------------------------------------------------------------------------
    # Generate Button
    # -------------------------------------------------------------------------
    col_btn, col_spacer = st.columns([1, 3])
    with col_btn:
        generate_clicked = st.button(
            label=get_text("generate_btn", lang),
            type="primary",
            use_container_width=True
        )

    st.markdown("---")

    # -------------------------------------------------------------------------
    # Pipeline Execution
    # -------------------------------------------------------------------------
    if generate_clicked:
        # Validation
        if not raw_text.strip():
            st.error(get_text("error_no_text", lang))
            st.stop()

        if not gemini_key.strip():
            st.error(get_text("error_no_gemini_key", lang))
            st.stop()

        if not openai_key.strip():
            st.error(get_text("error_no_openai_key", lang))
            st.stop()

        try:
            pptx_buffer = run_generation_pipeline(
                raw_text=raw_text,
                gemini_key=gemini_key,
                openai_key=openai_key,
                lang=lang,
                theme_name=selected_theme,
                num_slides=num_slides,
                progress_bar=progress_bar,
                status_text=status_text
            )

            # Success feedback
            st.balloons()
            st.success(get_text("success_msg", lang))

            # Download button
            st.download_button(
                label=get_text("download_btn", lang),
                data=pptx_buffer,
                file_name=f"Epic_Presentation_{selected_theme.replace(' ', '_')}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary",
                use_container_width=True
            )

        except RuntimeError as e:
            st.error(get_text("error_gemini", lang).format(e))
        except Exception as e:
            st.error(get_text("error_generic", lang).format(e))

    # -------------------------------------------------------------------------
    # Footer
    # -------------------------------------------------------------------------
    st.markdown(f'<div class="footer">{get_text("footer", lang)}</div>', unsafe_allow_html=True)


if __name__ == "__main__":
    main()
